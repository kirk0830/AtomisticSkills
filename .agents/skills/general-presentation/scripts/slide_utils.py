"""
Helper library for generating PowerPoint presentations from simulation results.

Provides high-level functions to build slides programmatically.
The agent writes standalone Python scripts that import from this module,
enabling iterative refinement of presentation content and layout.

Usage:
    import sys
    sys.path.insert(0, ".agents/skills/general-presentation/scripts")
    from slide_utils import *

Requirements:
    - Pixi environment: base
    - Required packages: python-pptx, Pillow
"""

from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image


# ---------------------------------------------------------------------------
# Theme configuration — override these before calling any builder functions
# ---------------------------------------------------------------------------

THEME = {
    # Colors
    "primary": RGBColor(0x1B, 0x3A, 0x5C),  # dark navy
    "secondary": RGBColor(0x2E, 0x86, 0xAB),  # teal accent
    "accent": RGBColor(0xA2, 0x3B, 0x72),  # magenta accent
    "background": RGBColor(0xFF, 0xFF, 0xFF),  # white
    "text_dark": RGBColor(0x2D, 0x2D, 0x2D),  # near-black
    "text_light": RGBColor(0xFF, 0xFF, 0xFF),  # white
    "text_muted": RGBColor(0x88, 0x88, 0x88),  # grey
    "header_bg": RGBColor(0x1B, 0x3A, 0x5C),  # same as primary
    "table_header_bg": RGBColor(0x1B, 0x3A, 0x5C),
    "table_row_alt": RGBColor(0xF2, 0xF6, 0xFA),  # light blue-grey
    "bullet_color": RGBColor(0x2E, 0x86, 0xAB),
    # Fonts
    "font_family": "Calibri",
    "title_size": Pt(32),
    "subtitle_size": Pt(20),
    "heading_size": Pt(28),
    "body_size": Pt(20),
    "caption_size": Pt(14),
    "table_size": Pt(14),
    # Slide dimensions (default 16:9)
    "width": Inches(13.333),
    "height": Inches(7.5),
    # Layout margins
    "margin_left": Inches(0.6),
    "margin_right": Inches(0.6),
    "margin_top": Inches(0.4),
    "header_height": Inches(0.9),
}


# ---------------------------------------------------------------------------
# Presentation lifecycle
# ---------------------------------------------------------------------------


def create_presentation(
    title: str = "Presentation",
    subtitle: str = "",
    author: str = "",
) -> Presentation:
    """Create a new presentation with a title slide.

    Args:
        title: Presentation title displayed on the first slide.
        subtitle: Subtitle text below the title.
        author: Author name shown on the title slide.

    Returns:
        A python-pptx Presentation object.
    """
    prs = Presentation()
    prs.slide_width = THEME["width"]
    prs.slide_height = THEME["height"]

    # Use blank layout for full control
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)

    # Full-slide background fill
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = THEME["primary"]

    # Title text
    _add_textbox(
        slide,
        left=THEME["margin_left"],
        top=Inches(2.2),
        width=prs.slide_width - THEME["margin_left"] - THEME["margin_right"],
        height=Inches(1.2),
        text=title,
        font_size=Pt(36),
        font_color=THEME["text_light"],
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )

    # Subtitle
    if subtitle:
        _add_textbox(
            slide,
            left=THEME["margin_left"],
            top=Inches(3.5),
            width=prs.slide_width - THEME["margin_left"] - THEME["margin_right"],
            height=Inches(0.6),
            text=subtitle,
            font_size=THEME["subtitle_size"],
            font_color=RGBColor(0xBB, 0xCC, 0xDD),
            alignment=PP_ALIGN.CENTER,
        )

    # Author
    if author:
        _add_textbox(
            slide,
            left=THEME["margin_left"],
            top=Inches(4.3),
            width=prs.slide_width - THEME["margin_left"] - THEME["margin_right"],
            height=Inches(0.4),
            text=author,
            font_size=THEME["caption_size"],
            font_color=RGBColor(0x99, 0xAA, 0xBB),
            alignment=PP_ALIGN.CENTER,
        )

    # Accent line
    _add_line(
        slide,
        left=Inches(4.0),
        top=Inches(3.3),
        width=prs.slide_width - Inches(8.0),
        color=THEME["secondary"],
        weight=Pt(2),
    )

    return prs


def save_presentation(prs: Presentation, path: str) -> str:
    """Save the presentation to disk.

    Args:
        prs: The Presentation object to save.
        path: Output file path (.pptx).

    Returns:
        The absolute path of the saved file.
    """
    out = Path(path).resolve()
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(f"Presentation saved: {out}")
    return str(out)


# ---------------------------------------------------------------------------
# Slide builders — each returns the Slide object for further customization
# ---------------------------------------------------------------------------


def add_title_slide(
    prs: Presentation,
    title: str,
    subtitle: str = "",
) -> "pptx.slide.Slide":
    """Add a title slide (useful for additional title pages mid-deck).

    Args:
        prs: The Presentation object.
        title: Slide title.
        subtitle: Optional subtitle.

    Returns:
        The created Slide object.
    """
    slide = _new_slide(prs)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = THEME["primary"]

    _add_textbox(
        slide,
        left=THEME["margin_left"],
        top=Inches(2.5),
        width=prs.slide_width - THEME["margin_left"] - THEME["margin_right"],
        height=Inches(1.0),
        text=title,
        font_size=Pt(32),
        font_color=THEME["text_light"],
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    if subtitle:
        _add_textbox(
            slide,
            left=THEME["margin_left"],
            top=Inches(3.6),
            width=prs.slide_width - THEME["margin_left"] - THEME["margin_right"],
            height=Inches(0.5),
            text=subtitle,
            font_size=THEME["subtitle_size"],
            font_color=RGBColor(0xBB, 0xCC, 0xDD),
            alignment=PP_ALIGN.CENTER,
        )
    return slide


def add_section_slide(
    prs: Presentation,
    title: str,
) -> "pptx.slide.Slide":
    """Add a section divider slide.

    Args:
        prs: The Presentation object.
        title: Section heading.

    Returns:
        The created Slide object.
    """
    slide = _new_slide(prs)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = THEME["secondary"]

    _add_textbox(
        slide,
        left=THEME["margin_left"],
        top=Inches(2.8),
        width=prs.slide_width - THEME["margin_left"] - THEME["margin_right"],
        height=Inches(1.0),
        text=title,
        font_size=Pt(36),
        font_color=THEME["text_light"],
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    return slide


def add_image_slide(
    prs: Presentation,
    title: str,
    image_path: str,
    caption: Optional[str] = None,
    notes: Optional[str] = None,
) -> "pptx.slide.Slide":
    """Add a slide with a single image (plot, structure visualization, etc.).

    The image is auto-scaled to fit the slide while maintaining aspect ratio.

    Args:
        prs: The Presentation object.
        title: Slide title.
        image_path: Path to the image file (PNG, JPG, etc.).
        caption: Optional caption below the image.
        notes: Optional speaker notes.

    Returns:
        The created Slide object.
    """
    slide = _new_slide(prs)
    _add_header(slide, prs, title)

    # Calculate image placement
    content_top = THEME["margin_top"] + THEME["header_height"] + Inches(0.2)
    content_width = prs.slide_width - THEME["margin_left"] - THEME["margin_right"]
    caption_reserve = Inches(0.5) if caption else Inches(0)
    max_img_height = prs.slide_height - content_top - Inches(0.4) - caption_reserve

    img_left, img_top, img_width, img_height = _fit_image(
        image_path,
        max_width=content_width,
        max_height=max_img_height,
    )
    # Center horizontally
    img_left = THEME["margin_left"] + (content_width - img_width) // 2
    img_top = content_top

    slide.shapes.add_picture(image_path, img_left, img_top, img_width, img_height)

    if caption:
        _add_textbox(
            slide,
            left=THEME["margin_left"],
            top=img_top + img_height + Inches(0.1),
            width=content_width,
            height=Inches(0.4),
            text=caption,
            font_size=THEME["caption_size"],
            font_color=THEME["text_muted"],
            alignment=PP_ALIGN.CENTER,
            italic=True,
        )

    if notes:
        slide.notes_slide.notes_text_frame.text = notes

    return slide


def add_two_image_slide(
    prs: Presentation,
    title: str,
    image_path_left: str,
    image_path_right: str,
    caption_left: Optional[str] = None,
    caption_right: Optional[str] = None,
) -> "pptx.slide.Slide":
    """Add a slide with two images side-by-side.

    Args:
        prs: The Presentation object.
        title: Slide title.
        image_path_left: Path to the left image.
        image_path_right: Path to the right image.
        caption_left: Optional caption for the left image.
        caption_right: Optional caption for the right image.

    Returns:
        The created Slide object.
    """
    slide = _new_slide(prs)
    _add_header(slide, prs, title)

    content_top = THEME["margin_top"] + THEME["header_height"] + Inches(0.2)
    total_width = prs.slide_width - THEME["margin_left"] - THEME["margin_right"]
    half_width = (total_width - Inches(0.4)) // 2  # gap between images
    caption_reserve = Inches(0.5) if (caption_left or caption_right) else Inches(0)
    max_img_height = prs.slide_height - content_top - Inches(0.4) - caption_reserve

    for i, (img_path, caption) in enumerate(
        [
            (image_path_left, caption_left),
            (image_path_right, caption_right),
        ]
    ):
        offset = (
            THEME["margin_left"]
            if i == 0
            else THEME["margin_left"] + half_width + Inches(0.4)
        )
        _, _, img_w, img_h = _fit_image(
            img_path, max_width=half_width, max_height=max_img_height
        )
        img_left = offset + (half_width - img_w) // 2
        slide.shapes.add_picture(img_path, img_left, content_top, img_w, img_h)

        if caption:
            _add_textbox(
                slide,
                left=offset,
                top=content_top + img_h + Inches(0.05),
                width=half_width,
                height=Inches(0.35),
                text=caption,
                font_size=THEME["caption_size"],
                font_color=THEME["text_muted"],
                alignment=PP_ALIGN.CENTER,
                italic=True,
            )

    return slide


def add_table_slide(
    prs: Presentation,
    title: str,
    headers: list[str],
    rows: list[list[str]],
) -> "pptx.slide.Slide":
    """Add a slide with a styled data table.

    Args:
        prs: The Presentation object.
        title: Slide title.
        headers: List of column header strings.
        rows: List of rows, each a list of cell strings.

    Returns:
        The created Slide object.
    """
    slide = _new_slide(prs)
    _add_header(slide, prs, title)

    content_top = THEME["margin_top"] + THEME["header_height"] + Inches(0.3)
    content_width = prs.slide_width - THEME["margin_left"] - THEME["margin_right"]
    num_rows = len(rows) + 1  # +1 for header
    num_cols = len(headers)

    row_height = Inches(0.4)
    table_height = row_height * num_rows
    max_table_height = prs.slide_height - content_top - Inches(0.5)
    if table_height > max_table_height:
        row_height = max_table_height // num_rows
        table_height = row_height * num_rows

    # Center table horizontally if it's narrower than content area
    table_width = min(content_width, Inches(2.0) * num_cols)
    table_left = THEME["margin_left"] + (content_width - table_width) // 2

    shape = slide.shapes.add_table(
        num_rows, num_cols, table_left, content_top, table_width, table_height
    )
    table = shape.table

    # Style header row
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = str(header)
        _style_cell(
            cell,
            bold=True,
            font_color=THEME["text_light"],
            fill_color=THEME["table_header_bg"],
        )

    # Style data rows
    for row_idx, row_data in enumerate(rows):
        fill_color = THEME["table_row_alt"] if row_idx % 2 == 1 else None
        for col_idx, value in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(value)
            _style_cell(cell, font_color=THEME["text_dark"], fill_color=fill_color)

    return slide


def add_bullets_slide(
    prs: Presentation,
    title: str,
    bullets: list[str],
) -> "pptx.slide.Slide":
    """Add a slide with bullet points.

    Args:
        prs: The Presentation object.
        title: Slide title.
        bullets: List of bullet-point strings.

    Returns:
        The created Slide object.
    """
    slide = _new_slide(prs)
    _add_header(slide, prs, title)

    content_top = THEME["margin_top"] + THEME["header_height"] + Inches(0.3)
    content_width = prs.slide_width - THEME["margin_left"] - THEME["margin_right"]
    content_height = prs.slide_height - content_top - Inches(0.4)

    txBox = slide.shapes.add_textbox(
        THEME["margin_left"] + Inches(0.3),
        content_top,
        content_width - Inches(0.6),
        content_height,
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet_text in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet_text
        p.font.size = THEME["body_size"]
        p.font.color.rgb = THEME["text_dark"]
        p.font.name = THEME["font_family"]
        p.space_after = Pt(14)
        p.level = 0
        # Bullet character
        p.bullet = True

    return slide


def add_image_and_text_slide(
    prs: Presentation,
    title: str,
    image_path: str,
    text_content: str,
    image_on_left: bool = True,
    caption: Optional[str] = None,
) -> "pptx.slide.Slide":
    """Add a slide with image on one side and text on the other.

    Args:
        prs: The Presentation object.
        title: Slide title.
        image_path: Path to the image file.
        text_content: Text to display beside the image.
        image_on_left: If True, image on left and text on right. Otherwise reversed.
        caption: Optional caption below the image.

    Returns:
        The created Slide object.
    """
    slide = _new_slide(prs)
    _add_header(slide, prs, title)

    content_top = THEME["margin_top"] + THEME["header_height"] + Inches(0.2)
    total_width = prs.slide_width - THEME["margin_left"] - THEME["margin_right"]
    half_width = (total_width - Inches(0.4)) // 2
    content_height = prs.slide_height - content_top - Inches(0.5)

    # Image side
    img_left = (
        THEME["margin_left"]
        if image_on_left
        else THEME["margin_left"] + half_width + Inches(0.4)
    )
    caption_reserve = Inches(0.4) if caption else Inches(0)
    _, _, img_w, img_h = _fit_image(
        image_path, max_width=half_width, max_height=content_height - caption_reserve
    )
    centered_left = img_left + (half_width - img_w) // 2
    slide.shapes.add_picture(image_path, centered_left, content_top, img_w, img_h)

    if caption:
        _add_textbox(
            slide,
            left=img_left,
            top=content_top + img_h + Inches(0.05),
            width=half_width,
            height=Inches(0.35),
            text=caption,
            font_size=THEME["caption_size"],
            font_color=THEME["text_muted"],
            alignment=PP_ALIGN.CENTER,
            italic=True,
        )

    # Text side
    text_left = (
        THEME["margin_left"] + half_width + Inches(0.4)
        if image_on_left
        else THEME["margin_left"]
    )
    _add_textbox(
        slide,
        left=text_left,
        top=content_top,
        width=half_width,
        height=content_height,
        text=text_content,
        font_size=THEME["body_size"],
        font_color=THEME["text_dark"],
    )

    return slide


# ---------------------------------------------------------------------------
# Internal helpers
# ---------------------------------------------------------------------------


def _new_slide(prs: Presentation) -> "pptx.slide.Slide":
    """Add a blank slide to the presentation."""
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def _add_header(
    slide: "pptx.slide.Slide",
    prs: Presentation,
    title: str,
) -> None:
    """Add the standard header bar with title to a content slide."""
    header_top = THEME["margin_top"]
    header_width = prs.slide_width - THEME["margin_left"] - THEME["margin_right"]

    # Header background rectangle
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        THEME["margin_left"],
        header_top,
        header_width,
        THEME["header_height"],
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = THEME["header_bg"]
    shape.line.fill.background()

    # Title text over the header
    _add_textbox(
        slide,
        left=THEME["margin_left"] + Inches(0.3),
        top=header_top,
        width=header_width - Inches(0.6),
        height=THEME["header_height"],
        text=title,
        font_size=THEME["heading_size"],
        font_color=THEME["text_light"],
        bold=True,
        vertical_anchor=MSO_ANCHOR.MIDDLE,
    )


def _add_textbox(
    slide: "pptx.slide.Slide",
    left: int,
    top: int,
    width: int,
    height: int,
    text: str,
    font_size: int = Pt(14),
    font_color: RGBColor = RGBColor(0x2D, 0x2D, 0x2D),
    bold: bool = False,
    italic: bool = False,
    alignment: PP_ALIGN = PP_ALIGN.LEFT,
    vertical_anchor: Optional[MSO_ANCHOR] = None,
) -> "pptx.shapes.placeholder.SlidePlaceholder":
    """Add a text box to a slide with standard formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    if vertical_anchor is not None:
        tf.paragraphs[0].alignment = alignment
        txBox.text_frame.auto_size = None
        # Set vertical centering via XML (python-pptx exposes this)
        from pptx.oxml.ns import qn

        txBody = txBox.text_frame._txBody
        bodyPr = txBody.find(qn("a:bodyPr"))
        anchor_map = {
            MSO_ANCHOR.TOP: "t",
            MSO_ANCHOR.MIDDLE: "ctr",
            MSO_ANCHOR.BOTTOM: "b",
        }
        bodyPr.set("anchor", anchor_map.get(vertical_anchor, "t"))

    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = font_color
    p.font.name = THEME["font_family"]
    p.font.bold = bold
    p.font.italic = italic
    p.alignment = alignment

    return txBox


def _add_line(
    slide: "pptx.slide.Slide",
    left: int,
    top: int,
    width: int,
    color: RGBColor,
    weight: int = Pt(1),
) -> None:
    """Add a horizontal line to a slide."""
    connector = slide.shapes.add_connector(
        1,  # straight connector
        left,
        top,
        left + width,
        top,
    )
    connector.line.color.rgb = color
    connector.line.width = weight


def _fit_image(
    image_path: str,
    max_width: int,
    max_height: int,
) -> tuple[int, int, int, int]:
    """Calculate image dimensions to fit within max bounds while preserving aspect ratio.

    Args:
        image_path: Path to the image file.
        max_width: Maximum width in EMU.
        max_height: Maximum height in EMU.

    Returns:
        Tuple of (left, top, width, height) in EMU. left and top are 0.
    """
    with Image.open(image_path) as img:
        img_w, img_h = img.size

    # Convert pixel dimensions to EMU for ratio calculation
    aspect = img_w / img_h

    # Try fitting to width first
    width = max_width
    height = int(width / aspect)

    # If too tall, fit to height instead
    if height > max_height:
        height = max_height
        width = int(height * aspect)

    return (0, 0, width, height)


def _style_cell(
    cell: "pptx.table._Cell",
    bold: bool = False,
    font_color: Optional[RGBColor] = None,
    fill_color: Optional[RGBColor] = None,
) -> None:
    """Apply styling to a table cell."""
    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.size = THEME["table_size"]
        paragraph.font.name = THEME["font_family"]
        if font_color:
            paragraph.font.color.rgb = font_color
        paragraph.font.bold = bold

    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Inches(0.08)
    cell.margin_right = Inches(0.08)
    cell.margin_top = Inches(0.04)
    cell.margin_bottom = Inches(0.04)

    if fill_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill_color


def get_approx_text_width_inches(text: str, font_size: int = 10) -> float:
    """Estimate the physical width of text in inches.

    Args:
        text: The string to measure.
        font_size: Font size in points.

    Returns:
        Estimated width in inches.
    """
    base_char_width = 0.08 * (font_size / 10.0)
    padding = 0.15
    return (len(text) * base_char_width) + padding


def add_autofit_box(
    slide: "pptx.slide.Slide",
    left: int,
    top: int,
    height: int,
    text: str,
    bg_color: RGBColor,
    font_color: RGBColor = RGBColor(0xFF, 0xFF, 0xFF),
    font_size: int = 10,
    is_rounded: bool = True,
    bold: bool = False,
) -> tuple["pptx.shapes.autoshape.Shape", float]:
    """Add a colored box that automatically scales its width to fit the text.

    Args:
        slide: The Slide object.
        left: Left position in EMU.
        top: Top position in EMU.
        height: Box height in EMU.
        text: Text content.
        bg_color: Background color.
        font_color: Text color.
        font_size: Font size in points.
        is_rounded: If True, uses rounded rectangle.
        bold: If True, makes text bold.

    Returns:
        A tuple of (Shape object, calculated width in inches).
    """
    width_inches = get_approx_text_width_inches(text, font_size)
    width = Inches(width_inches)

    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if is_rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = bg_color

    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.alignment = PP_ALIGN.CENTER
    shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    return shape, width_inches
